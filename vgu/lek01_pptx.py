from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = 'Привет!'
subtitle.text = 'Я -- презентация!'
prs.save('test.pptx')

# создаем новый слайд со схемой для добавления изображений
slide = prs.slides.add_slide(prs.slide_layouts[8])
slide.shapes.title.text = "А теперь с картинкой"

# Добавляем изображение
print("Контейнер для изображения:", slide.placeholders[1])
placeholder = slide.placeholders[1]
placeholder.insert_picture('pic.jpeg')
prs.save('test.pptx')
